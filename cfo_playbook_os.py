import os
import sys
import time
import io
from datetime import datetime

# --- NOTE ---
# Versione Clean: Rimosso il sistema di auto-installazione instabile.
# Le dipendenze vengono gestite esternamente (es. tramite ripara_installazione.bat)

# --- IMPORTAZIONI PROTETTE ---
try:
    import streamlit as st
    import pandas as pd
    import numpy as np
    import plotly.graph_objects as go
    import numpy_financial as npf
    from pptx import Presentation
    from fpdf import FPDF
except ImportError as e:
    # Se manca una libreria, ora vediamo ESATTAMENTE quale manca
    try:
        import streamlit as st
        st.error(f"❌ ERRORE MANCANZA LIBRERIA: {e}")
        st.warning("Per risolvere: Chiudi il programma ed esegui il file 'ripara_installazione.bat'.")
        st.stop()
    except ImportError:
        # Se manca proprio Streamlit, lo stampiamo nella finestra nera
        print(f"!!! ERRORE CRITICO: Manca la libreria '{e.name}' !!!")
        print("Esegui il file 'ripara_installazione.bat' per installarla.")
        time.sleep(10)
        sys.exit()

# --- TAVOLO COLORI EXECUTIVE ---
COLORS = {
    "primary": "#1a1a1a",      # Nero Black Swan
    "accent": "#deff9a",       # Lime/Oro Infografica
    "secondary": "#3B82F6",    # Blu Professionale
    "danger": "#EF4444",       # Rosso Rischio
    "success": "#10B981",      # Verde Profitto
    "neutral": "#64748b"       # Grigio
}

# --- DIZIONARIO LINGUE ESTESO ---
LANGUAGES = {
    "Italiano": {
        "sidebar_title": "CFO OS 2025",
        "settings": "Configurazione Generale",
        "company": "Azienda",
        "currency": "Valuta",
        "tabs": ["🏠 Sintesi", "📖 Guida CEO", "📊 Investimenti", "🔄 SaaS", "💧 Liquidità", "🎯 Break-even", "🌪️ Stress Test"],
        "labels": {
            "settings": "Configurazione", "company": "Azienda", "currency": "Valuta",
            "capex": "Investimento (Capex)", "years": "Anni Orizzonte", "wacc": "WACC %",
            "rev1": "Ricavi Anno 1", "growth": "Crescita Ricavi %", "cogs": "COGS %",
            "opex": "OPEX Anno 1", "opex_g": "Crescita OPEX %", "tax": "Tax Rate %",
            "arr": "ARR Iniziale", "exp": "Espansione €", "churn_val": "Churn €",
            "cac": "CAC", "arpu": "ARPU Mensile", "churn_rate": "Churn Mensile %", "gm_saas": "Gross Margin SaaS %",
            "cash": "Cassa", "debt_st": "Debiti Breve T.", "debt_lt": "Debiti Lungo T.",
            "dso": "DSO (Incasso)", "dio": "DIO (Magazzino)", "dpo": "DPO (Fornitori)",
            "price": "Prezzo Unitario", "var_cost": "Costo Variabile Unit.", "fix_cost": "Costi Fissi Tot.", "vol": "Volume",
            "shock_rev": "Shock Ricavi %", "shock_cost": "Shock Costi %",
            "gen_report": "Genera Report", "download": "Scarica"
        },
        "tips": {
            "wacc": "Il 'Costo del Capitale'. Rappresenta il rendimento minimo che devi ottenere per soddisfare banche e azionisti.",
            "capex": "Soldi che escono SUBITO. Include macchinari, software, ristrutturazioni.",
            "cogs": "Costi diretti (materie prime). Se aumentano, il margine lordo scende.",
            "opex": "Spese fisse (affitti, stipendi amministrativi).",
            "nopat": "Net Operating Profit After Tax. È il vero utile operativo 'pulito' dagli interessi sul debito.",
            "npv": "La somma di tutti i soldi futuri portati al valore di oggi. Se è > 0, procedi.",
            "tax": "Aliquota fiscale stimata. Default 28%."
        },
        "kpi": {
            "npv": "NPV (VAN)", "irr": "IRR (TIR)", "payback": "Payback", "pfn": "PFN",
            "nrr": "Retention (NRR)", "ltv_cac": "LTV/CAC", "ccc": "Ciclo Cassa (CCC)", "safety": "Margine Sicurezza",
            "bep": "BEP (Valore)"
        },
        "recom": {
            "npv_ok": "✅ Semaforo Verde: Il progetto crea valore economico reale.",
            "npv_ko": "❌ Semaforo Rosso: Il progetto distrugge ricchezza. Non approvare.",
            "liq_ok": "💧 Cassa OK: Nessuna tensione di liquidità a breve.",
            "liq_ko": "⚠️ Allerta Cassa: Ciclo troppo lungo, rischi di finire i soldi.",
            "saas_ok": "🚀 Motore SaaS Sano: Ottima efficienza commerciale.",
            "saas_ko": "🔻 Motore SaaS Rotto: Spendi troppo per acquisire clienti.",
            "stress_ok": "🛡️ Resiliente: L'azienda regge lo shock sui ricavi.",
            "stress_ko": "🌪️ Fragile: Lo stress test porta l'EBITDA in negativo."
        },
        "titles": {
            "invest": "Analisi Investimenti", "saas": "Metriche SaaS", "liq": "Liquidità & PFN",
            "bep": "Break-even Point", "stress": "Stress Test", "sum": "Executive Summary"
        },
        "f_ricavi": "Ricavi", "f_ebitda": "EBITDA", "f_nopat": "NOPAT",
        "headers": {
            "bep_intro": "Analisi del punto di pareggio", 
            "stress_intro": "Simulazione scenari di crisi", 
            "dash_ceo": "Dashboard Direzionale",
            "recom_strat": "Raccomandazioni Strategiche",
            "details": "📋 Dettaglio Annuale",
            "cf_chart": "📉 Dinamica dei Flussi di Cassa (Waterfall)",
            "drivers": "Driver Operativi",
            "spread_chart": "⚖️ Economic Value Spread (ROIC vs WACC)"
        },
        "guide": {
            "title": "Manuale Strategico per il CEO",
            "faq_title": "❓ 5. FAQ",
            "intro": """
            **1. INTRODUZIONE**
            
            Benvenuto nel **CFO Operating System 2025**, uno strumento decisionale strategico progettato per i leader finanziari moderni.
            A differenza dei fogli di calcolo tradizionali, questo sistema ti consente di simulare scenari, testare la resilienza del tuo modello di business e convalidare le decisioni di investimento in tempo reale.

            **Filosofia Core:**
            *"La finanza non riguarda il reporting del passato. Riguarda l'architettura del valore futuro."*
            
            ---
            **2. PER INIZIARE**
            
            * **Installazione:** Assicurati che Python sia installato. Esegui il file `.bat` o lo script da terminale.
            * **Configurazione (Sidebar):** Inserisci Nome Azienda, Valuta (€, $, £) e Lingua.
            """,
            "modules_title": "3. PANORAMICA MODULI",
            "mod_invest": """
            ### 📊 Tab 1: Analisi Investimenti (Capex)
            **Obiettivo:** Decidere se un nuovo progetto/investimento è finanziariamente sostenibile.
            
            **Input:**
            * *Capex:* L'esborso di cassa iniziale (es. macchinari, R&D, software).
            * *Orizzonte:* Quanti anni vuoi proiettare.
            * *WACC:* Il tuo costo del capitale (benchmark: 8-12%).
            * *Driver Ricavi:* Ricavi attesi Anno 1 e % di crescita annuale.
            * *Margini:* % Target EBITDA e struttura OPEX.

            **Metriche Chiave:**
            * **NPV (Valore Attuale Netto):** La metrica più importante.
                * *Positivo (+):* Il progetto crea ricchezza. **APPROVA**.
                * *Negativo (-):* Il progetto distrugge valore. **RIFIUTA**.
            * **IRR (Tasso Interno di Rendimento):** Il rendimento annualizzato del progetto. Deve essere > WACC.
            * **Payback Period:** Tempo necessario per recuperare l'esborso iniziale.
            """,
            "mod_saas": """
            ### 🔄 Tab 2: Metriche SaaS & Subscription
            **Obiettivo:** Valutare la salute dei modelli di ricavo ricorrenti.
            
            **Input:**
            * *ARR:* Annual Recurring Revenue.
            * *Churn:* % di ricavi persi annualmente.
            * *CAC:* Costo per acquisire un singolo cliente (Marketing + Vendite).
            * *ARPU:* Average Revenue Per User.

            **Metriche Chiave:**
            * **LTV/CAC Ratio:** La "Metrica Aurea" dell'economia unitaria.
                * *> 3.0x:* Crescita sana.
                * *< 1.0x:* Perdi soldi su ogni cliente.
            * **NRR (Net Revenue Retention):** Misura la crescita dai clienti esistenti. Target > 100%.
            """,
            "mod_liq": """
            ### 💧 Tab 3: Liquidità & Debito Netto
            **Obiettivo:** Assicurare che l'azienda abbia abbastanza cassa per sopravvivere e crescere.
            
            **Input:**
            * *Cassa Disponibile:* Saldo bancario attuale.
            * *Debito:* Prestiti a breve e lungo termine.
            * *Capitale Circolante:* DSO (Vendite), DIO (Scorte), DPO (Fornitori).

            **Metriche Chiave:**
            * **PFN (Posizione Finanziaria Netta):** Debito Totale meno Cassa.
            * **CCC (Cash Conversion Cycle):** Il numero di giorni in cui la cassa è "bloccata" nelle operazioni. Più basso è, meglio è.
            """,
            "mod_bep": """
            ### 🎯 Tab 4: Analisi Break-even
            **Obiettivo:** Definire la soglia minima di viabilità.
            
            **Input:**
            * *Prezzo Unitario & Costo Variabile:* Per calcolare il Margine di Contribuzione.
            * *Costi Fissi:* Affitto, stipendi, spese generali.

            **Metriche Chiave:**
            * **BEP (Valore):** L'importo esatto di ricavi necessario per raggiungere profitto zero.
            * **Margine di Sicurezza:** La % di cui i ricavi possono scendere prima di iniziare a perdere soldi (Target > 20%).
            """,
            "mod_stress": """
            ### 🌪️ Tab 5: Stress Test (Black Swan)
            **Obiettivo:** Testare la resilienza contro shock di mercato.
            
            **Azione:** Muovi i cursori per simulare un evento "Cigno Nero" (es. -20% Ricavi, +15% Costi).
            **Risultato:** L'EBITDA rimane positivo? Se sì, il tuo modello di business è **Anti-fragile**.
            
            ---
            **4. EXECUTIVE SUMMARY & EXPORT**
            Il Tab **Sintesi** aggrega i dati per fornire lo stato in tempo reale (Semafori Verdi/Rossi) e Raccomandazioni Strategiche.
            Usa i pulsanti nella sidebar per scaricare report PDF e PPT.
            """,
            "faq_q1": "**Q: Perché calcolate il NOPAT?**",
            "faq_a1": "A: L'EBITDA è spesso fuorviante. Il NOPAT (Net Operating Profit After Tax) è il vero flusso di cassa operativo disponibile per ripagare gli investitori.",
            "faq_q2": "**Q: Cosa succede se il Payback è 'N.D.'?**",
            "faq_a2": "A: Significa che il progetto non raggiunge mai il pareggio entro il periodo selezionato. È un investimento altamente rischioso.",
            "faq_q3": "**Q: Posso usarlo per aziende non-SaaS?**",
            "faq_a3": "A: Sì. Salta semplicemente il tab 'SaaS'. I moduli Investimenti, Liquidità e Stress Test sono universali."
        },
        "footer_base": "Basato su The Black Swan CFO Playbook"
    },
    "English": {
        "sidebar_title": "CFO OS 2025",
        "settings": "General Settings",
        "company": "Company",
        "currency": "Currency",
        "tabs": ["🏠 Summary", "📖 CEO Guide", "📊 Investments", "🔄 SaaS", "💧 Liquidity", "🎯 Break-even", "🌪️ Stress Test"],
        "labels": {
            "settings": "Configuration", "company": "Company", "currency": "Currency",
            "capex": "Initial Investment", "years": "Time Horizon", "wacc": "WACC %",
            "rev1": "Year 1 Revenue", "growth": "Revenue Growth %", "cogs": "COGS %",
            "opex": "Year 1 OPEX", "opex_g": "OPEX Growth %", "tax": "Tax Rate %",
            "arr": "Starting ARR", "exp": "Expansion €", "churn_val": "Churn €",
            "cac": "CAC", "arpu": "Monthly ARPU", "churn_rate": "Monthly Churn %", "gm_saas": "SaaS Gross Margin %",
            "cash": "Cash on Hand", "debt_st": "Short Term Debt", "debt_lt": "Long Term Debt",
            "dso": "DSO (Receivables)", "dio": "DIO (Inventory)", "dpo": "DPO (Payables)",
            "price": "Unit Price", "var_cost": "Unit Var. Cost", "fix_cost": "Total Fixed Costs", "vol": "Volume",
            "shock_rev": "Revenue Shock %", "shock_cost": "Cost Shock %",
            "gen_report": "Generate Report", "download": "Download"
        },
        "tips": {
            "wacc": "Weighted Average Cost of Capital. Minimum return required.",
            "capex": "Immediate cash outflow for fixed asset acquisition.",
            "cogs": "Direct costs related to production or service.",
            "opex": "Fixed operating costs (rent, salaries, admin).",
            "nopat": "Net Operating Profit After Tax.",
            "npv": "Net Present Value: wealth created today by future flows.",
            "tax": "Estimated average tax rate."
        },
        "kpi": {
            "npv": "NPV", "irr": "IRR", "payback": "Payback", "pfn": "Net Debt",
            "nrr": "Retention (NRR)", "ltv_cac": "LTV/CAC", "ccc": "Cash Cycle (CCC)", "safety": "Safety Margin",
            "bep": "BEP (Value)"
        },
        "recom": {
            "npv_ok": "✅ Green Light: Project creates real value.",
            "npv_ko": "❌ Red Light: Project destroys wealth.",
            "liq_ok": "💧 Liquidity OK: Efficient cash cycle.",
            "liq_ko": "⚠️ Cash Alert: Cycle too long, risk of crisis.",
            "saas_ok": "🚀 Healthy SaaS Engine: Great sales efficiency.",
            "saas_ko": "🔻 Broken SaaS Engine: Acquisition cost too high.",
            "stress_ok": "🛡️ Resilient: Company withstands revenue shock.",
            "stress_ko": "🌪️ Fragile: Stress test pushes EBITDA to negative."
        },
        "titles": {
            "invest": "Investment Analysis", "saas": "SaaS Metrics", "liq": "Liquidity & Net Debt",
            "bep": "Break-even Point", "stress": "Stress Test", "sum": "Executive Summary"
        },
        "f_ricavi": "Revenue", "f_ebitda": "EBITDA", "f_nopat": "NOPAT",
        "headers": {
            "bep_intro": "Break-even analysis", 
            "stress_intro": "Crisis simulation scenarios", 
            "dash_ceo": "Executive Dashboard",
            "recom_strat": "Strategic Recommendations",
            "details": "📋 Annual Detail",
            "cf_chart": "📉 Cash Flow Dynamics (Waterfall)",
            "drivers": "Operational Drivers",
            "spread_chart": "⚖️ Economic Value Spread (ROIC vs WACC)"
        },
        "guide": {
            "title": "CEO Strategic Manual",
            "faq_title": "❓ 5. FAQ",
            "intro": """
            **1. INTRODUCTION**
            
            Welcome to the **CFO Operating System 2025**, a strategic decision-making tool designed for modern financial leaders.
            Unlike traditional spreadsheets, this system allows you to simulate scenarios, stress-test your business model, and validate investment decisions in real-time.

            **Core Philosophy:**
            *"Finance is not about reporting the past. It is about architectural design of future value."*
            
            ---
            **2. GETTING STARTED**
            
            * **Installation:** Ensure Python is installed. Run the `.bat` file or script via terminal.
            * **Configuration (Sidebar):** Enter Company Name, Currency (€, $, £), and Language.
            """,
            "modules_title": "3. MODULES OVERVIEW",
            "mod_invest": """
            ### 📊 Tab 1: Investment Analysis (Capex)
            **Goal:** Decide if a new project/investment is financially viable.
            
            **Inputs:**
            * *Capex:* The initial cash outlay (e.g., machinery, R&D, software).
            * *Horizon:* How many years you want to project.
            * *WACC:* Your cost of capital (benchmark: 8-12%).
            * *Revenue Drivers:* Expected Year 1 revenue and annual growth %.
            * *Margins:* Target EBITDA % and OPEX structure.

            **Key Metrics:**
            * **NPV (Net Present Value):** The single most important metric. 
                * *Positive (+):* The project creates wealth. **APPROVE**.
                * *Negative (-):* The project destroys value. **REJECT**.
            * **IRR (Internal Rate of Return):** The annualized return of the project. Must be > WACC.
            * **Payback Period:** Time required to recover the initial cash outlay.
            """,
            "mod_saas": """
            ### 🔄 Tab 2: SaaS & Subscription Metrics
            **Goal:** Evaluate the health of recurring revenue models.
            
            **Inputs:**
            * *ARR:* Annual Recurring Revenue.
            * *Churn:* % of revenue lost annually.
            * *CAC:* Cost to acquire a single customer (Marketing + Sales).
            * *ARPU:* Average Revenue Per User.

            **Key Metrics:**
            * **LTV/CAC Ratio:** The "Golden Metric" of unit economics.
                * *> 3.0x:* Healthy growth.
                * *< 1.0x:* You are losing money on every customer.
            * **NRR (Net Revenue Retention):** Measures growth from existing customers. Target > 100%.
            """,
            "mod_liq": """
            ### 💧 Tab 3: Liquidity & Net Debt
            **Goal:** Ensure the company has enough cash to survive and grow.
            
            **Inputs:**
            * *Cash on Hand:* Current bank balance.
            * *Debt:* Short-term and Long-term loans.
            * *Working Capital Days:* DSO (Sales), DIO (Inventory), DPO (Payables).

            **Key Metrics:**
            * **PFN (Net Financial Position):** Total Debt minus Cash.
            * **CCC (Cash Conversion Cycle):** The number of days your cash is "trapped" in operations. Lower is better.
            """,
            "mod_bep": """
            ### 🎯 Tab 4: Break-even Analysis
            **Goal:** Define the minimum viability threshold.
            
            **Inputs:**
            * *Unit Price & Variable Cost:* To calculate the Contribution Margin.
            * *Fixed Costs:* Rent, salaries, overheads.

            **Key Metrics:**
            * **BEP (Value):** The exact revenue amount needed to reach zero profit.
            * **Safety Margin:** The % your revenue can drop before you start losing money. (Target > 20%).
            """,
            "mod_stress": """
            ### 🌪️ Tab 5: Stress Test (Black Swan)
            **Goal:** Test resilience against market shocks.
            
            **Action:** Move the sliders to simulate a "Black Swan" event (e.g., -20% Revenue, +15% Costs).
            **Result:** Does EBITDA remain positive? If yes, your business model is **Anti-fragile**.
            
            ---
            **4. EXECUTIVE SUMMARY & EXPORT**
            The **Home Tab (Summary)** acts as your cockpit. It aggregates data from all other modules to provide real-time status and recommendations.
            Use the buttons in the sidebar to download PDF and PPT reports.
            """,
            "faq_q1": "**Q: Why do you calculate NOPAT?**",
            "faq_a1": "A: EBITDA is often misleading. NOPAT (Net Operating Profit After Tax) is the true operating cash flow available to pay back investors.",
            "faq_q2": "**Q: What if Payback is 'N.D.'?**",
            "faq_a2": "A: It means the project never breaks even within the selected timeframe. It is a highly risky investment.",
            "faq_q3": "**Q: Can I use this for non-SaaS companies?**",
            "faq_a3": "A: Yes. Just skip the 'SaaS' tab. The Investment, Liquidity, and Stress Test modules are universal."
        },
        "footer_base": "Based on The Black Swan CFO Playbook"
    }
}

# --- FUNZIONI DI EXPORT ---
def sanitize_text(text):
    if not isinstance(text, str): text = str(text)
    return text.replace('€', 'EUR').replace('£', 'GBP').replace('$', 'USD').encode('latin-1', 'replace').decode('latin-1')

def generate_pdf(data, recoms, title):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 10, sanitize_text(title), ln=True, align='C')
    pdf.ln(10)
    pdf.set_font("Helvetica", "B", 12)
    pdf.cell(0, 10, "KPI Summary:", ln=True)
    pdf.set_font("Helvetica", "", 11)
    for k, v in data.items():
        pdf.cell(0, 8, f"{sanitize_text(k)}: {sanitize_text(v)}", ln=True)
    pdf.ln(10)
    pdf.set_font("Helvetica", "B", 12)
    pdf.cell(0, 10, "Strategic Recommendations:", ln=True)
    pdf.set_font("Helvetica", "", 10)
    width = pdf.w - 2 * pdf.l_margin
    for r in recoms:
        pdf.set_x(pdf.l_margin)
        pdf.multi_cell(width, 8, sanitize_text(f"- {r}"))
    pdf.ln(10)
    pdf.set_font("Helvetica", "I", 8)
    pdf.cell(0, 10, f"Generated by Black Swan CFO OS - {datetime.now().strftime('%Y-%m-%d')}", ln=True, align='C')
    return bytes(pdf.output())

def generate_pptx(data, recoms, title):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"Analisi Finanziaria Strategica\n{datetime.now().strftime('%Y-%m-%d')}"
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Performance Indicators"
    tf = slide.placeholders[1].text_frame
    for k, v in data.items():
        p = tf.add_paragraph()
        p.text = f"{k}: {v}"
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Strategic Recommendations"
    tf = slide.placeholders[1].text_frame
    for r in recoms:
        p = tf.add_paragraph()
        p.text = r
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

def generate_csv(data):
    df = pd.DataFrame(list(data.items()), columns=['Metric', 'Value'])
    return df.to_csv(index=False).encode('utf-8')

# --- CONFIGURAZIONE PAGINA STREAMLIT ---
st.set_page_config(page_title="Black Swan CFO OS", layout="wide", page_icon="🦢")

st.markdown(f"""
<style>
    .logo-container {{background-color: {COLORS['primary']}; padding: 20px; border-radius: 10px; border: 2px solid {COLORS['accent']}; text-align: center; margin-bottom: 20px;}}
    .logo-title {{color: {COLORS['accent']}; font-family: 'Arial', sans-serif; font-weight: bold; font-size: 28px; margin: 0; letter-spacing: 2px;}}
    .logo-subtitle {{color: #FFFFFF; font-family: sans-serif; font-size: 12px; letter-spacing: 4px; margin-top: 5px;}}
    .stMetric {{background-color: white; border: 1px solid #e0e0e0; border-radius: 8px; padding: 15px;}}
</style>
""", unsafe_allow_html=True)

# --- SIDEBAR ---
with st.sidebar:
    st.markdown("""<div class="logo-container"><div class="logo-title">BLACK SWAN</div><div class="logo-subtitle">CFO PLAYBOOK</div></div>""", unsafe_allow_html=True)
    lang_key = st.selectbox("Language / Lingua", list(LANGUAGES.keys()))
    L = LANGUAGES[lang_key]
    Labels = L["labels"]
    Tips = L["tips"]
    
    st.divider()
    st.subheader(Labels["settings"])
    azienda = st.text_input(Labels["company"], "Alpha Industries Inc.")
    valuta = st.selectbox(Labels["currency"], ["€", "$", "£"])

# --- TAB RENDERING ---
tabs = st.tabs(L["tabs"])

# ================= TAB 1: GUIDA CEO =================
with tabs[1]:
    st.header(L["guide"]["title"])
    st.markdown(L["guide"]["intro"])
    
    with st.expander(L["guide"]["modules_title"], expanded=True):
        st.markdown(L["guide"]["mod_invest"])
        st.markdown(L["guide"]["mod_saas"])
        st.markdown(L["guide"]["mod_liq"])
        st.markdown(L["guide"]["mod_bep"])
        st.markdown(L["guide"]["mod_stress"])
    
    with st.expander(L["guide"]["faq_title"]):
        st.markdown(f"{L['guide']['faq_q1']}\n{L['guide']['faq_a1']}")
        st.markdown(f"{L['guide']['faq_q2']}\n{L['guide']['faq_a2']}")
        st.markdown(f"{L['guide']['faq_q3']}\n{L['guide']['faq_a3']}")

# ================= TAB 2: INVESTIMENTI =================
with tabs[2]:
    st.header(L["titles"]["invest"])
    
    with st.expander(Labels["settings"], expanded=True):
        c1, c2, c3 = st.columns(3)
        inv = c1.number_input(Labels["capex"], value=500000, help=Tips["capex"])
        durata = c2.slider(Labels["years"], 1, 15, 5)
        wacc = c3.slider(Labels["wacc"], 1.0, 20.0, 10.0, help=Tips["wacc"]) / 100
        
    st.subheader(L["headers"]["drivers"])
    c4, c5, c6 = st.columns(3)
    rev1 = c4.number_input(Labels["rev1"], value=300000)
    growth = c5.slider(Labels["growth"], -10.0, 50.0, 15.0) / 100
    cogs_p = c6.slider(Labels["cogs"], 0, 90, 40, help=Tips["cogs"]) / 100
    
    c7, c8, c9 = st.columns(3)
    opex1 = c7.number_input(Labels["opex"], value=50000, help=Tips["opex"])
    opex_g = c8.slider(Labels["opex_g"], 0.0, 20.0, 3.0) / 100
    tax_r = c9.slider(Labels["tax"], 0, 50, 28, help=Tips["tax"]) / 100
    
    # --- CALCOLI ---
    cf_list = [-inv]
    roic_list = []
    years_labels = ["Anno 0 (CAPEX)"] + [f"Anno {i}" for i in range(1, durata + 1)]
    da_annuo = inv / durata if durata > 0 else 0
    rows = []
    
    for i in range(1, durata + 1):
        r_t = rev1 * ((1 + growth) ** (i-1))
        c_t = r_t * cogs_p
        o_t = opex1 * ((1 + opex_g) ** (i-1))
        ebitda_t = r_t - c_t - o_t
        ebit_t = ebitda_t - da_annuo
        tax_val = max(0, ebit_t * tax_r)
        nopat_t = ebit_t - tax_val
        fcf_t = nopat_t + da_annuo
        
        cf_list.append(fcf_t)
        # Rendimento calcolato sul CAPEX iniziale per spread
        roic_t = (nopat_t / inv) * 100 if inv > 0 else 0
        roic_list.append(roic_t)
        
        rows.append({
            "Anno": i, 
            L["f_ricavi"]: r_t, 
            L["f_ebitda"]: ebitda_t, 
            L["f_nopat"]: nopat_t,
            "FCF": fcf_t
        })
        
    npv_val = npf.npv(wacc, cf_list)
    irr_val = npf.irr(cf_list)
    cum_fcf = np.cumsum(cf_list)
    
    # KPI Principali
    m1, m2, m3 = st.columns(3)
    m1.metric(L["kpi"]["npv"], f"{valuta} {npv_val:,.0f}", delta="OK" if npv_val > 0 else "KO")
    m2.metric(L["kpi"]["irr"], f"{irr_val:.1%}" if irr_val else "N/A")
    m3.metric(L["kpi"]["payback"], f"{next((i for i, v in enumerate(cum_fcf) if v >= 0), 'N.D.')} Anni")

    # --- GRAFICO 1: WATERFALL PAYBACK ---
    st.subheader(L["headers"]["cf_chart"])
    fig_cf = go.Figure()
    fig_cf.add_trace(go.Bar(
        x=years_labels, y=cf_list, 
        name="FCF Annuo", 
        marker=dict(color=[COLORS['danger'] if x < 0 else COLORS['secondary'] for x in cf_list]), 
        text=[f"{val/1000:.0f}k" for val in cf_list],
        textposition='outside'
    ))
    fig_cf.add_trace(go.Scatter(
        x=years_labels, y=cum_fcf, 
        name="Flusso Cumulato", 
        line=dict(color=COLORS['accent'], width=4),
        mode='lines+markers',
        marker=dict(size=10, line=dict(color="white", width=2))
    ))
    fig_cf.update_layout(
        template="plotly_white", 
        hovermode="x unified",
        xaxis_title="Orizzonte Temporale",
        yaxis_title=f"Cash Flow ({valuta})",
        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1)
    )
    st.plotly_chart(fig_cf, use_container_width=True)

    # --- GRAFICO 2: VALUE SPREAD ---
    st.subheader(L["headers"]["spread_chart"])
    fig_spread = go.Figure()
    fig_spread.add_trace(go.Bar(
        x=years_labels[1:], y=roic_list, 
        name="ROIC (Rendimento)", 
        marker=dict(color=COLORS['accent'])
    ))
    fig_spread.add_trace(go.Scatter(
        x=years_labels[1:], y=[wacc*100]*durata, 
        name="WACC (Costo Hurdle)", 
        line=dict(color=COLORS['danger'], width=3, dash='dash'),
        mode='lines'
    ))
    fig_spread.update_layout(yaxis_ticksuffix="%", template="plotly_white", xaxis_title="Anni")
    st.plotly_chart(fig_spread, use_container_width=True)

    # Tabella Dettaglio
    st.subheader(L["headers"]["details"])
    st.dataframe(pd.DataFrame(rows).style.format("{:,.0f}"), use_container_width=True)

# ================= TAB 3: SAAS =================
with tabs[3]:
    st.header(L["titles"]["saas"])
    s1, s2 = st.columns(2)
    arr = s1.number_input(Labels["arr"], value=1000000)
    churn = s2.slider(Labels["churn_rate"], 0.1, 10.0, 2.0) / 100
    arpu = st.number_input(Labels["arpu"], value=500)
    cac = st.number_input(Labels["cac"], value=4000)
    
    ltv = (arpu * 0.85) / churn if churn > 0 else 0
    ratio = ltv / cac if cac > 0 else 0
    st.metric(L["kpi"]["ltv_cac"], f"{ratio:.1f}x", delta="Eccellente" if ratio > 3 else "Critico")

# ================= TAB 4: LIQUIDITA =================
with tabs[4]:
    st.header(L["titles"]["liq"])
    l1, l2 = st.columns(2)
    cash = l1.number_input(Labels["cash"], value=150000)
    debt = l2.number_input(Labels["debt_lt"], value=400000)
    pfn = debt - cash
    st.metric(L["kpi"]["pfn"], f"{valuta} {pfn:,.0f}", delta_color="inverse")
    
    st.divider()
    d1, d2, d3 = st.columns(3)
    dso = d1.number_input(Labels["dso"], value=60)
    dio = d2.number_input(Labels["dio"], value=45)
    dpo = d3.number_input(Labels["dpo"], value=90)
    ccc = dso + dio - dpo
    
    # --- GRAFICO CCC ORIZZONTALE ---
    st.subheader(f"Composizione {L['kpi']['ccc']}")
    fig_ccc = go.Figure()
    fig_ccc.add_trace(go.Bar(y=["Ciclo"], x=[dso], name="DSO (Incasso)", orientation='h', marker=dict(color=COLORS['secondary'])))
    fig_ccc.add_trace(go.Bar(y=["Ciclo"], x=[dio], name="DIO (Magazzino)", orientation='h', marker=dict(color=COLORS['neutral'])))
    fig_ccc.add_trace(go.Bar(y=["Ciclo"], x=[-dpo], name="DPO (Debiti)", orientation='h', marker=dict(color=COLORS['accent'])))
    fig_ccc.update_layout(barmode='relative', template="plotly_white", height=250, title=f"Ciclo Totale: {ccc} giorni")
    st.plotly_chart(fig_ccc, use_container_width=True)

# ================= TAB 5: BREAK-EVEN =================
with tabs[5]:
    st.header(L["titles"]["bep"])
    st.write(L["headers"]["bep_intro"])
    b1, b2 = st.columns(2)
    price = b1.number_input(Labels["price"], value=100)
    vc = b1.number_input(Labels["var_cost"], value=60)
    fc = b1.number_input(Labels["fix_cost"], value=150000)
    vol = b2.number_input(Labels["vol"], value=5000)
    
    mc = price - vc
    bep_val = (fc / mc) * price if mc > 0 else 0
    safety = ((vol * price) - bep_val) / (vol * price) if vol > 0 else 0
    
    st.metric(L["kpi"]["bep"], f"{valuta} {bep_val:,.0f}")
    st.metric(L["kpi"]["safety"], f"{safety:.1%}")

# ================= TAB 6: STRESS TEST =================
with tabs[6]:
    st.header(L["titles"]["stress"])
    shock = st.slider(Labels["shock_rev"], -50, 0, -20)
    
    base_eb = rev1 * (1 - cogs_p) - opex1
    stress_eb = (rev1 * (1 + shock/100)) * (1 - cogs_p) - opex1
    
    # --- GRAFICO STRESS COMPARISON ---
    fig_stress = go.Figure()
    fig_stress.add_trace(go.Bar(
        x=["Scenario Base", "Scenario Shock"], 
        y=[base_eb, stress_eb], 
        marker=dict(color=[COLORS['primary'], COLORS['danger']]),
        text=[f"{val/1000:.0f}k" for val in [base_eb, stress_eb]],
        textposition='auto'
    ))
    fig_stress.update_layout(title="Impatto EBITDA sullo Scenario Black Swan", template="plotly_white")
    st.plotly_chart(fig_stress, use_container_width=True)

# ================= TAB 0: SINTESI & EXPORT =================
with tabs[0]:
    st.header(L["titles"]["sum"])
    k1, k2, k3, k4 = st.columns(4)
    k1.metric(L["kpi"]["npv"], f"{valuta} {npv_val:,.0f}")
    k2.metric(L["kpi"]["pfn"], f"{valuta} {pfn:,.0f}")
    k3.metric(L["kpi"]["ltv_cac"], f"{ratio:.1f}x")
    k4.metric(L["kpi"]["safety"], f"{safety:.1%}")
    
    recoms = []
    if npv_val > 0: recoms.append(L["recom"]["npv_ok"])
    else: recoms.append(L["recom"]["npv_ko"])
    if ccc < 60: recoms.append(L["recom"]["liq_ok"])
    else: recoms.append(L["recom"]["liq_ko"])
    if ratio > 3: recoms.append(L["recom"]["saas_ok"])
    else: recoms.append(L["recom"]["saas_ko"])
    if stress_eb > 0: recoms.append(L["recom"]["stress_ok"])
    else: recoms.append(L["recom"]["stress_ko"])
    
    st.subheader(L["headers"]["recom_strat"])
    for r in recoms:
        if any(icon in r for icon in ["✅", "🚀", "🛡️", "💧"]):
            st.success(r)
        else:
            st.error(r)
        
    # Preparazione Dati Export
    export_data = {
        L["kpi"]["npv"]: f"{valuta} {npv_val:,.0f}",
        L["kpi"]["irr"]: f"{irr_val:.1%}" if irr_val else "N/A",
        L["kpi"]["pfn"]: f"{valuta} {pfn:,.0f}",
        L["kpi"]["ltv_cac"]: f"{ratio:.2f}x",
        L["kpi"]["safety"]: f"{safety:.1%}",
        "Data Estrazione": datetime.now().strftime("%Y-%m-%d")
    }
    
    st.divider()
    st.subheader(Labels["gen_report"])
    c_pdf, c_ppt, c_csv = st.columns(3)
    
    with c_pdf:
        st.download_button("📕 PDF", data=generate_pdf(export_data, recoms, f"{azienda} - Strategic Report"), file_name=f"{azienda}_Report.pdf")
    with c_ppt:
        st.download_button("📙 PPTX", data=generate_pptx(export_data, recoms, f"{azienda} Executive"), file_name=f"{azienda}_Executive.pptx")
    with c_csv:
        st.download_button("📗 CSV", data=generate_csv(export_data), file_name=f"{azienda}_Dati.csv")

st.divider()
st.caption(f"Black Swan CFO OS v11.0 | {azienda} | {L['footer_base']}")